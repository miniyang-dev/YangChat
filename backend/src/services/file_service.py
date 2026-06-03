"""
file_service.py — 解析上傳檔案為純文字，並提供圖片 vision 描述（透過 Gemini Flash）
支援格式：PDF (.pdf), PowerPoint (.pptx), Word (.docx), 純文字 (.txt), Markdown (.md), CSV (.csv)
圖片格式：JPEG (.jpg/.jpeg), PNG (.png), GIF (.gif), WebP (.webp)
"""
import base64
import csv
import io
from pathlib import Path

import fitz  # pymupdf
import httpx
from docx import Document
from pptx import Presentation


MAX_FILE_SIZE_MB = 10
MAX_CHARS = 20000  # 避免塞爆 context window

SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".txt", ".docx", ".md", ".csv"}


def extract_text(filename: str, content: bytes) -> str:
    """
    根據副檔名選擇對應解析器，回傳純文字（最多 MAX_CHARS 字元）。
    不支援的格式拋出 ValueError。
    """
    suffix = Path(filename).suffix.lower()

    if len(content) > MAX_FILE_SIZE_MB * 1024 * 1024:
        raise ValueError(f"檔案超過 {MAX_FILE_SIZE_MB}MB 限制")

    if suffix in (".txt", ".md"):
        text = _extract_txt(content)
    elif suffix == ".pdf":
        text = _extract_pdf(content)
    elif suffix == ".pptx":
        text = _extract_pptx(content)
    elif suffix == ".docx":
        text = _extract_docx(content)
    elif suffix == ".csv":
        text = _extract_csv(content)
    else:
        supported = "、".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"不支援的檔案格式：{suffix}，請上傳 {supported}")

    return text[:MAX_CHARS]


def _extract_txt(content: bytes) -> str:
    """嘗試 UTF-8，失敗則 fallback 到 big5"""
    try:
        return content.decode("utf-8")
    except UnicodeDecodeError:
        return content.decode("big5", errors="replace")


def _extract_pdf(content: bytes) -> str:
    """用 pymupdf 逐頁提取文字"""
    doc = fitz.open(stream=content, filetype="pdf")
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text().strip()
        if text:
            pages.append(f"[第 {i + 1} 頁]\n{text}")
    doc.close()
    return "\n\n".join(pages)


def _extract_pptx(content: bytes) -> str:
    """用 python-pptx 逐張投影片提取文字"""
    prs = Presentation(io.BytesIO(content))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides.append(f"[投影片 {i + 1}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_docx(content: bytes) -> str:
    """用 python-docx 提取段落文字"""
    doc = Document(io.BytesIO(content))
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            paragraphs.append(text)
    # 也提取表格內容
    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                paragraphs.append(" | ".join(row_texts))
    return "\n".join(paragraphs)


def _extract_csv(content: bytes) -> str:
    """解析 CSV，轉成可讀的純文字格式（標頭 + 資料列）"""
    try:
        text = content.decode("utf-8")
    except UnicodeDecodeError:
        text = content.decode("big5", errors="replace")

    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return ""

    lines = []
    headers = rows[0]
    lines.append("欄位：" + "、".join(headers))
    lines.append(f"共 {len(rows) - 1} 筆資料：")

    for i, row in enumerate(rows[1:], 1):
        if len(row) == len(headers):
            pairs = [f"{h}={v}" for h, v in zip(headers, row) if v]
            lines.append(f"[第{i}筆] " + "，".join(pairs))
        else:
            lines.append(f"[第{i}筆] " + "，".join(row))

    return "\n".join(lines)


# ── Image Vision（Gemini Flash）─────────────────────────────────────────────

SUPPORTED_IMAGE_TYPES = {
    "image/jpeg": "jpeg",
    "image/png":  "png",
    "image/gif":  "gif",
    "image/webp": "webp",
}

MAX_IMAGE_SIZE = 10 * 1024 * 1024  # 10MB


async def describe_image(image_bytes: bytes, content_type: str, api_key: str) -> str:
    """
    把圖片送給 Gemini Flash，回傳圖片的詳細文字描述。
    描述結果會當成 file_context 注入給 Pioneer AI。
    """
    if content_type not in SUPPORTED_IMAGE_TYPES:
        raise ValueError(f"不支援的圖片格式：{content_type}，支援 JPEG / PNG / GIF / WebP")

    if len(image_bytes) > MAX_IMAGE_SIZE:
        raise ValueError("圖片大小超過 10MB 限制")

    # Base64 編碼
    b64 = base64.b64encode(image_bytes).decode("utf-8")
    mime = content_type  # e.g. "image/jpeg"

    payload = {
        "contents": [{
            "parts": [
                {
                    "inline_data": {
                        "mime_type": mime,
                        "data": b64,
                    }
                },
                {
                    "text": (
                        "請詳細描述這張圖片的所有內容，包括：\n"
                        "1. 圖片中的主要元素、物件、人物\n"
                        "2. 文字內容（如有，請完整抄錄）\n"
                        "3. 圖表、數據或結構性資訊（如有）\n"
                        "4. 整體場景與背景\n"
                        "5. 任何可能有助於分析或回答問題的細節\n"
                        "請用繁體中文回答，盡量完整詳細。"
                    )
                }
            ]
        }],
        "generationConfig": {
            "maxOutputTokens": 2048,
            "temperature": 0.1,
        }
    }

    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent"

    async with httpx.AsyncClient(timeout=30) as client:
        # W-2: API key 用 header 傳送，避免出現在 access log / URL 中
        resp = await client.post(url, json=payload, headers={"x-goog-api-key": api_key})
        resp.raise_for_status()
        data = resp.json()

    # 取出描述文字
    try:
        description = data["candidates"][0]["content"]["parts"][0]["text"]
    except (KeyError, IndexError) as e:
        raise ValueError(f"Gemini 回應格式異常：{e}，原始回應：{str(data)[:200]}")

    return description.strip()
