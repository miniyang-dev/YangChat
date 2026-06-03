"""
export_service.py — 將 AI 產出的投影片結構組裝成 .pptx
"""
import io
from typing import List

from pptx import Presentation
from pptx.util import Inches, Pt


def build_pptx(slides_data: List[dict]) -> bytes:
    """
    輸入格式：
    [
        {"title": "投影片標題", "content": ["重點1", "重點2", ...]},
        ...
    ]
    回傳：.pptx 的 bytes
    """
    prs = Presentation()
    # title_only_layout（index 5）：有標題佔位符，可自由加文字框
    title_only_layout = prs.slide_layouts[5]

    for slide_data in slides_data:
        slide = prs.slides.add_slide(title_only_layout)

        # ── 標題 ──
        title_box = slide.shapes.title
        title_box.text = slide_data.get("title", "")
        title_para = title_box.text_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True

        # ── 內容文字框 ──
        content_lines = slide_data.get("content", [])
        if content_lines:
            txBox = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(9), Inches(5)
            )
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, line in enumerate(content_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {line}"
                p.font.size = Pt(18)

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
