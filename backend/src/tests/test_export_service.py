"""
test_export_service.py — 測試 PPTX 產出服務
"""
import io
import sys
import os
import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from services.export_service import build_pptx
from pptx import Presentation


def test_build_pptx_returns_bytes():
    slides = [{"title": "測試", "content": ["重點A", "重點B"]}]
    result = build_pptx(slides)
    assert isinstance(result, bytes)
    assert len(result) > 0


def test_build_pptx_correct_slide_count():
    slides = [
        {"title": "第一頁", "content": ["A"]},
        {"title": "第二頁", "content": ["B"]},
        {"title": "第三頁", "content": ["C"]},
    ]
    result = build_pptx(slides)
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 3


def test_build_pptx_empty_slides():
    result = build_pptx([])
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 0


def test_build_pptx_slide_title_content():
    slides = [{"title": "AI 趨勢", "content": ["GPT-4", "Gemini", "Claude"]}]
    result = build_pptx(slides)
    prs = Presentation(io.BytesIO(result))
    slide = prs.slides[0]
    # 檢查標題
    assert slide.shapes.title.text == "AI 趨勢"
    # 檢查內容文字框有 bullet points
    all_text = " ".join(
        shape.text_frame.text
        for shape in slide.shapes
        if shape.has_text_frame and shape != slide.shapes.title
    )
    assert "GPT-4" in all_text
    assert "Gemini" in all_text


def test_build_pptx_no_content():
    """只有標題，沒有 content 的投影片"""
    slides = [{"title": "空白投影片", "content": []}]
    result = build_pptx(slides)
    prs = Presentation(io.BytesIO(result))
    assert len(prs.slides) == 1
    assert prs.slides[0].shapes.title.text == "空白投影片"
